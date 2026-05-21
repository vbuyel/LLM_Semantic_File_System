import io
import logging
import os
import re
import unicodedata
from collections.abc import Callable
from pathlib import Path
from typing import Optional

from adapters.kafka import KafkaOperations
from domain.domain import SendToKafka

logger = logging.getLogger(__name__)

_ENCODINGS = ("utf-8", "latin-1", "cp1252")

import pymupdf
from docx import Document
from pptx import Presentation
import openpyxl


class TextExtractorOperations:
    MAX_CHUNK_CHARS = 150 * 1024  # 150KB per Kafka message — keeps each msg safely under broker limits
    _TEXT_EXTENSIONS = {
       # Documentation / markup
        ".txt", ".md", ".rst", ".rtf", ".tex", ".org", ".adoc", ".wiki",
        ".log", ".csv", ".tsv",
        # Web
        ".html", ".htm", ".xhtml", ".css", ".scss", ".less", ".sass",
        ".xml", ".json", ".yaml", ".yml", ".toml",
        # Config / infra
        ".ini", ".cfg", ".conf", ".env", ".gitignore", ".editorconfig",
        ".dockerfile", ".makefile", ".cmake", ".gradle",
        # Source code
        ".py", ".pyw", ".js", ".mjs", ".cjs", ".ts", ".jsx", ".tsx",
        ".java", ".c", ".cpp", ".cc", ".cxx", ".h", ".hpp", ".hxx",
        ".rs", ".go", ".rb", ".php", ".phtml", ".swift", ".kt", ".kts", ".scala",
        ".sh", ".bash", ".zsh", ".fish", ".ps1", ".bat", ".cmd",
        ".pl", ".pm", ".lua", ".r", ".m", ".mm",
        ".sql", ".graphql", ".gql", ".proto",
        ".svelte", ".vue",
    }
    _DOC_EXTENSIONS = {".pdf", ".docx", ".pptx", ".xlsx"}
    _DEPENDENCY_ERRORS: dict[str, str] = {
        ".pdf": "PyMuPDF required for PDF: pip install pymupdf",
        ".docx": "python-docx required for DOCX: pip install python-docx",
        ".pptx": "python-pptx required for PPTX: pip install python-pptx",
        ".xlsx": "openpyxl required for XLSX: pip install openpyxl",
    }
    _SUPPORTED_EXTENSIONS = _TEXT_EXTENSIONS | _DOC_EXTENSIONS


    def __init__(self):
        self.kafka = KafkaOperations()


    def extract_text_from_bytes(self, content: bytes, ext: str) -> str:
        """Extract text from file content bytes, bypassing the filesystem.

        Returns an empty string for unsupported or binary file types.
        """
        ext = self._normalise_ext(ext)
        if ext not in self._SUPPORTED_EXTENSIONS:
            return ""

        return self._bytes_extractor_for(ext)(content)


    def extract_text_from_file(self, file_path: str) -> str:
        """Extract text from a local file path.

        Returns an empty string for unsupported or binary file types.
        """
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"File not found: {file_path}")

        ext = self._normalise_ext(Path(file_path).suffix)
        if ext not in self._SUPPORTED_EXTENSIONS:
            return ""

        return self._file_extractor_for(ext)(file_path)


    def clean_text(self, text: str) -> str:
        """Strip null bytes and control characters, normalise whitespace."""
        cleaned: list[str] = []
        for ch in text:
            cat = unicodedata.category(ch)
            if ch in ("\t", "\n", "\r"):
                cleaned.append(ch)
            elif cat[0] == "C":
                cleaned.append(" ")
            else:
                cleaned.append(ch)
        return re.compile(r"[^\S\t\n\r]+").sub(" ", "".join(cleaned)).strip()


    def is_readable(
        self, text: str, min_chars: int = 10, threshold: float = 0.3
    ) -> bool:
        """Return True when *text* has enough meaningful content to be worth indexing."""
        if not text or len(text) < min_chars:
            return False
        readable = sum(1 for ch in text if ch.isalnum() or ch in (" ", "\t", "\n"))
        return readable / len(text) >= threshold


    def chunk_text(self, text: str, max_chars: int = MAX_CHUNK_CHARS) -> list[str]:
        """Split text at word boundaries, keeping each chunk ≤ max_chars."""
        if len(text) <= max_chars:
            return [text]
        words = text.split()
        chunks: list[str] = []
        current: list[str] = []
        current_len = 0
        for word in words:
            word_len = len(word)
            sep = 1 if current else 0
            if current_len + sep + word_len > max_chars and current:
                chunks.append(" ".join(current))
                current = [word]
                current_len = word_len
            else:
                current_len += sep + word_len
                current.append(word)
        if current:
            chunks.append(" ".join(current))
        return chunks


    async def send_chunked_kafka(
        self,
        action: str,
        file_name: str,
        file_path: str,
        text: str,
        owner: Optional[str],
        storage_type: str,
        file_size: int = 0,
    ) -> None:
        """Clean text, split into readable chunks, and send each chunk to Kafka."""
        text = self.clean_text(text)
        if not self.is_readable(text):
            logger.info(f"Skipping {file_path}: no readable text after cleaning")
            return

        raw_chunks = self.chunk_text(text)
        chunks = [c for c in raw_chunks if self.is_readable(c)]
        if not chunks:
            logger.info(f"Skipping {file_path}: no readable chunks after filtering")
            return

        for i, chunk in enumerate(chunks):
            try:
                await self.kafka.send_command(
                    SendToKafka(
                        action=action,
                        file_name=file_name,
                        file_path=file_path,
                        text=chunk,
                        owner=owner,
                        storage_type=storage_type,
                        chunk_index=i,
                        file_size=file_size,
                    )
                )
            except Exception as e:
                logger.warning(
                    f"Failed to send Kafka event (chunk {i+1}/{len(chunks)}): {e}"
                )


    def _normalise_ext(self, ext: str) -> str:
        """Normalise a file extension."""
        ext = ext.strip().lower()
        if not ext.startswith("."):
            ext = "." + ext
        return ext


    def _bytes_extractor_for(self, ext: str) -> Callable[[bytes], str]:
        """Get a bytes extractor for a file extension."""
        return {
            ".pdf": self._extract_pdf_bytes,
            ".docx": self._extract_docx_bytes,
            ".pptx": self._extract_pptx_bytes,
            ".xlsx": self._extract_xlsx_bytes,
        }.get(ext, self._decode_bytes)


    def _file_extractor_for(self, ext: str) -> Callable[[str], str]:
        """Get a file extractor for a file extension."""
        return {
            ".pdf": self._extract_pdf,
            ".docx": self._extract_docx,
            ".pptx": self._extract_pptx,
            ".xlsx": self._extract_xlsx,
        }.get(ext, self._read_text_file)


    def _decode_bytes(self, content: bytes) -> str:
        """Decode bytes to a string."""
        for encoding in _ENCODINGS:
            try:
                return content.decode(encoding)
            except UnicodeDecodeError:
                continue
        return ""


    def _read_text_file(self, file_path: str) -> str:
        """Read a text file."""
        for encoding in _ENCODINGS:
            try:
                with open(file_path, "r", encoding=encoding) as f:
                    return f.read()
            except UnicodeDecodeError:
                continue
        return ""


    def _extract_pdf_bytes(self, content: bytes) -> str:
        """Extract text from PDF bytes."""
        doc = pymupdf.open(stream=content, filetype="pdf")
        try:
            return "\n".join(page.get_text() for page in doc)
        finally:
            doc.close()


    def _extract_pdf(self, file_path: str) -> str:
        """Extract text from PDF file."""
        doc = pymupdf.open(file_path)
        try:
            return "\n".join(page.get_text() for page in doc)
        finally:
            doc.close()


    def _extract_docx_bytes(self, content: bytes) -> str:
        """Extract text from DOCX bytes."""
        doc = Document(io.BytesIO(content))
        response = "\n".join(p.text for p in doc.paragraphs)
        return response


    def _extract_docx(self, file_path: str) -> str:
        """Extract text from DOCX file."""
        doc = Document(file_path)
        response = "\n".join(p.text for p in doc.paragraphs)
        return response


    def _collect_pptx_text(self, prs: Presentation) -> str:
        """Collect text from PPTX."""
        texts: list[str] = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texts.append(shape.text)
        response = "\n".join(texts)
        return response


    def _extract_pptx_bytes(self, content: bytes) -> str:
        """Extract text from PPTX bytes."""
        presentation = Presentation(io.BytesIO(content))
        response = self._collect_pptx_text(presentation)
        return response


    def _extract_pptx(self, file_path: str) -> str:
        """Extract text from PPTX file."""
        presentation = Presentation(file_path)
        response = self._collect_pptx_text(presentation)
        return response


    def _collect_xlsx_rows(self, wb: openpyxl.Workbook) -> str:
        """Collect rows from XLSX."""
        rows: list[str] = []
        try:
            for sheet in wb.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    cells = [str(c) for c in row if c is not None]
                    if cells:
                        rows.append("\t".join(cells))
        finally:
            wb.close()
        response = "\n".join(rows)
        return response


    def _extract_xlsx_bytes(self, content: bytes) -> str:
        """Extract text from XLSX bytes."""
        wb = openpyxl.load_workbook(io.BytesIO(content), read_only=True, data_only=True)
        response = self._collect_xlsx_rows(wb)
        return response


    def _extract_xlsx(self, file_path: str) -> str:
        """Extract text from XLSX file."""
        wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
        response = self._collect_xlsx_rows(wb)
        return response
